import os
import re
from PIL import Image 
from pytesseract import pytesseract 
from pptx import Presentation
import PyPDF2

# Read API key if needed for future audio/video features
try:
    with open("api_key.txt") as file:
        api_key = file.readline().strip()
except FileNotFoundError:
    api_key = None

class Transcriber:
    def __init__(self, file_path):
        self.file_path = file_path
        self.media_type = file_path.split(".")[-1].lower()
        print(f"Processing file type: {self.media_type}")
        
    def _clean_text(self, text):
        """Clean and normalize text output"""
        if isinstance(text, list):
            return [self._clean_text(item) for item in text]
            
        # Remove extra whitespace, normalize newlines, clean up special characters
        text = str(text).strip()
        text = re.sub(r'\n+', '\n', text)  # Replace multiple newlines with single
        text = re.sub(r'[ \t]+', ' ', text)  # Replace multiple spaces/tabs with single space
        text = re.sub(r'\s*\n\s*', '\n', text)  # Clean spaces around newlines
        text = re.sub(r'[^\x00-\x7F]+', ' ', text)  # Remove non-ASCII characters
        return text.strip()

    def image_transcribe(self):
        """Extract text from images using OCR"""
        try:
            img = Image.open(self.file_path) 
            text = pytesseract.image_to_string(img)
            return self._clean_text(text)
        except Exception as e:
            raise Exception(f"Error processing image: {str(e)}")
        finally:
            # Clean up temp file
            if os.path.exists(self.file_path):
                os.remove(self.file_path)

    def ppt_transcribe(self):
        """Extract text from PowerPoint presentations"""
        try:
            prs = Presentation(self.file_path)
            slide_texts = []
            
            for slide in prs.slides: 
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text.strip()
                        if text: 
                            slide_texts.append(text)
            
            return self._clean_text(slide_texts)  # Return as cleaned list
        except Exception as e:
            raise Exception(f"Error processing PowerPoint: {str(e)}")
        finally:
            # Clean up temp file
            if os.path.exists(self.file_path):
                os.remove(self.file_path)

    def pdf_transcribe(self):
        """Extract text from PDF files"""
        try:
            with open(self.file_path, "rb") as pdf_file:
                pdf_reader = PyPDF2.PdfReader(pdf_file)
                text_pages = []
                
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    page_text = page.extract_text()
                    if page_text.strip():  # Only add non-empty pages
                        text_pages.append(page_text)
                
                full_text = "\n".join(text_pages)
                return self._clean_text(full_text)
        except Exception as e:
            raise Exception(f"Error processing PDF: {str(e)}")
        finally:
            # Clean up temp file
            if os.path.exists(self.file_path):
                os.remove(self.file_path)

def runner(media, media_type):
    """
    Route to appropriate transcription method based on media type
    """
    try:
        if media_type[0] == 'application':
            if media_type[1] == 'pdf':
                return media.pdf_transcribe()
            elif 'presentation' in media_type[1]:  # PPT files
                return media.ppt_transcribe()
        elif media_type[0] == 'image':
            return media.image_transcribe()
        else:
            raise Exception(f"Unsupported media type: {'/'.join(media_type)}")
            
    except Exception as e:
        raise Exception(f"Error in transcription: {str(e)}")
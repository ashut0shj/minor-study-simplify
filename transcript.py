import moviepy.editor as mp
import assemblyai as aai
import os
from PIL import Image 
from pytesseract import pytesseract 
from pptx import Presentation
import PyPDF2

with open("api_key.txt") as file:
    file_content = file.readline()
    api_key = file_content  

class Transcriber:

    aai.settings.api_key = api_key

    def __init__(self,file_path):
        self.transcriber = aai.Transcriber()
        self.file_path = file_path
        self.media_type = file_path.split(".")[-1]
        print(self.media_type)
        
    
    def audio_transcribe(self):
        print("Transcribing")
        try:
            transcript = self.transcriber.transcribe(r"temp.mp3")
        except:
            transcript = self.transcriber.transcribe(r"temp.mpeg")
        os.remove(self.file_path)
        self.transcript = transcript.text
        return transcript.text

    def video_transcribe(self):
        video_clip = mp.VideoFileClip(self.file_path)
        video_clip.audio.write_audiofile('temp.wav')
        video_clip.close()
        os.remove('temp.mp4')
        print("video converted")
        return self.audio_transcribe()
    
    def image_transcribe(self):
        image_path = self.file_path
        img = Image.open(image_path) 
        text = pytesseract.image_to_string(img) 
        os.remove(self.file_path)
        self.transcript = (text[:-1])
        return text[:-1]


    def ppt_transcribe(self):
        prs = Presentation(self.file_path)
        os.remove('temp.vnd.openxmlformats-officedocument.presentationml.presentation')
        slide_titles = [] 
        for slide in prs.slides: 
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text.strip()
                    if text: 
                        slide_titles.append(text)
        self.transcript = slide_titles #data is in a list
        return self.transcript

    def pdf_transcribe(self):
        pdf_file = open(self.file_path, "rb")
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        text = []
        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            text.append(page.extract_text())
        pdf_file.close()
        full_text = "\n".join(text)
        self.transcript = full_text
        os.remove(self.file_path)
        return self.transcript
                
    def printt(self):
        print(self.transcript)
        
        
def runner(media, media_type):
    if media_type[0] == 'video':
            transcript = media.video_transcribe()
    elif media_type[0] == 'audio':
        transcript = media.audio_transcribe()
    elif media_type[0] == 'vnd.openxmlformats-officedocument.presentationml.presentation':
        transcript = media.ppt_transcribe()
    elif media_type[1] == 'pdf':
        transcript = media.pdf_transcribe()
    elif media_type[0] == 'image':
        transcript = media.image_transcribe()
    else:
        print("else stateent")
        transcript = media.ppt_transcribe()

    return transcript
from pptx import Presentation

prs = Presentation(r"C:\Users\ashut\Downloads\Minor Project 3rd Sem PPT.pptx") # load the ppt
slide_titles = [] # container for slide titles

for slide in prs.slides: # iterate over each slide
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text.strip()
            if text:  # Check if text is not empty
                slide_titles.append(text)

# Print the extracted slide titles
for title in slide_titles:
    print(title)
